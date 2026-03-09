"""
Extract all outfit items from PPTX slides for AI analysis pipeline.

Handles Canva exports where images are embedded as blip fills inside FREEFORM shapes.
Extracts panel colors from GROUP shapes.
Tracks image deduplication across looks.

Outputs:
  public/outfits/extracted/look-{N}/item-{n}.{ext}
  public/outfits/extracted/manifest-raw.json
"""

import hashlib
import json
import os
import re
import struct
import sys
from pptx import Presentation

PPTX_PATH = sys.argv[1] if len(sys.argv) > 1 else os.path.expanduser(
    "~/Downloads/melisaonder.com outfits (1).pptx"
)
OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "public", "outfits", "extracted")
MAX_SLIDES = int(sys.argv[2]) if len(sys.argv) > 2 else 0  # 0 = all

EMU_PER_PX = 914400 / 96


def emu_to_px(emu: int) -> int:
    return round(emu / EMU_PER_PX)


def get_image_dimensions(blob: bytes, content_type: str) -> tuple[int, int] | None:
    """Get actual pixel dimensions from image blob."""
    try:
        if content_type == "image/png" and len(blob) > 24:
            w = struct.unpack(">I", blob[16:20])[0]
            h = struct.unpack(">I", blob[20:24])[0]
            return (w, h)
        elif content_type in ("image/jpeg", "image/jpg") and len(blob) > 2:
            # Parse JPEG markers for SOF
            i = 2
            while i < len(blob) - 1:
                if blob[i] != 0xFF:
                    break
                marker = blob[i + 1]
                if marker in (0xC0, 0xC2):  # SOF0 or SOF2
                    h = struct.unpack(">H", blob[i + 5:i + 7])[0]
                    w = struct.unpack(">H", blob[i + 7:i + 9])[0]
                    return (w, h)
                length = struct.unpack(">H", blob[i + 2:i + 4])[0]
                i += 2 + length
    except Exception:
        pass
    return None


def blob_hash(blob: bytes) -> str:
    """Short hash for deduplication tracking."""
    return hashlib.md5(blob).hexdigest()[:12]


def extract_blip_image(shape, slide_part):
    """Extract image blob from a shape's blip fill (Canva freeform pattern)."""
    xml = shape._element.xml
    rids = re.findall(r'r:embed="(rId\d+)"', xml)
    for rid in rids:
        try:
            rel = slide_part.rels[rid]
            if "image" in rel.reltype:
                return rel.target_part.blob, rel.target_part.content_type
        except Exception:
            continue
    return None, None


def extract_panel_colors(shapes):
    """Extract background panel colors from GROUP shapes."""
    colors = []
    for shape in shapes:
        if shape.shape_type != 6:  # GROUP
            continue
        try:
            xml = shape._element.xml
            # Look for solid fill colors in the group's freeform children
            hex_matches = re.findall(r'<a:srgbClr val="([A-Fa-f0-9]{6})"', xml)
            for hex_color in hex_matches:
                color = f"#{hex_color}"
                if color not in colors:
                    colors.append(color)
        except Exception:
            pass
    return colors


def extract_items(shapes, slide_part, out_dir, global_hashes):
    """Extract image items from slide shapes."""
    items = []

    for shape_idx, shape in enumerate(shapes):
        # Recurse into groups for any nested images
        if shape.shape_type == 6:
            try:
                items.extend(extract_items(shape.shapes, slide_part, out_dir, global_hashes))
            except Exception:
                pass
            continue

        img_bytes = None
        content_type = None

        # Method 1: standard picture shape
        if hasattr(shape, "image"):
            try:
                img_bytes = shape.image.blob
                content_type = shape.image.content_type
            except Exception:
                pass

        # Method 2: Canva-style blip fill in freeform/auto shapes
        if img_bytes is None:
            img_bytes, content_type = extract_blip_image(shape, slide_part)

        if img_bytes is None:
            continue

        try:
            ext = content_type.split("/")[-1].replace("jpeg", "jpg")
            left_px = emu_to_px(shape.left)
            top_px = emu_to_px(shape.top)
            w_px = emu_to_px(shape.width)
            h_px = emu_to_px(shape.height)
            rotation = shape.rotation if hasattr(shape, "rotation") else 0

            # Get actual image dimensions
            actual_dims = get_image_dimensions(img_bytes, content_type)

            # Deduplication hash
            img_hash = blob_hash(img_bytes)

            # Save file
            fname = f"item-{len(items)}.{ext}"
            fpath = os.path.join(out_dir, fname)
            with open(fpath, "wb") as f:
                f.write(img_bytes)

            # Track global hash for cross-slide dedup
            if img_hash not in global_hashes:
                global_hashes[img_hash] = []

            item = {
                "file": fname,
                "left": left_px,
                "top": top_px,
                "width": w_px,
                "height": h_px,
                "rotation": rotation,
                "zIndex": len(items),
                "imageHash": img_hash,
                "fileSize": len(img_bytes),
                "contentType": content_type,
            }

            if actual_dims:
                item["actualWidth"] = actual_dims[0]
                item["actualHeight"] = actual_dims[1]

            items.append(item)
        except Exception as e:
            print(f"  Skipping shape {shape_idx}: {e}")

    return items


def slide_to_jpg_name(slide_idx: int, total_slides: int) -> str | None:
    """Map PPTX slide index to existing slide JPG filename."""
    # Slides 3-51 (idx 2-50) → casual-1 through casual-49
    if 2 <= slide_idx <= 50:
        return f"casual-{slide_idx - 1}.jpg"
    # Slide 52 (idx 51) is a divider
    if slide_idx == 51:
        return None
    # Slides 53-86 (idx 52-85) → formal-1 through formal-34
    if 52 <= slide_idx <= 85:
        return f"formal-{slide_idx - 51}.jpg"
    return None


def main():
    print(f"Opening: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)

    slide_w = emu_to_px(prs.slide_width)
    slide_h = emu_to_px(prs.slide_height)
    total_slides = len(list(prs.slides))
    print(f"Slide dimensions: {slide_w}x{slide_h} px")
    print(f"Total slides: {total_slides}")

    os.makedirs(OUT_DIR, exist_ok=True)

    manifest = {
        "slideWidth": slide_w,
        "slideHeight": slide_h,
        "totalSlides": total_slides,
        "looks": [],
    }

    # Global hash tracking for cross-slide deduplication
    global_hashes: dict[str, list[str]] = {}

    slides = list(prs.slides)
    if MAX_SLIDES > 0:
        slides = slides[:MAX_SLIDES]

    skipped = 0
    extracted = 0

    for slide_idx, slide in enumerate(slides):
        # Skip title/divider slides (indices 0, 1, 51)
        if slide_idx in (0, 1, 51):
            skipped += 1
            continue

        look_id = f"look-{slide_idx + 1}"
        slide_dir = os.path.join(OUT_DIR, look_id)
        os.makedirs(slide_dir, exist_ok=True)

        # Map to existing slide JPG
        jpg_name = slide_to_jpg_name(slide_idx, total_slides)
        slide_image = f"/outfits/slides/{jpg_name}" if jpg_name else None

        # Determine category from JPG name
        category = None
        if jpg_name:
            category = "casual" if jpg_name.startswith("casual") else "formal"

        print(f"Extracting slide {slide_idx + 1}/{len(slides)}: {look_id} ({jpg_name or 'no-jpg'})")

        # Extract panel colors
        panel_colors = extract_panel_colors(slide.shapes)

        # Extract items
        items = extract_items(slide.shapes, slide.part, slide_dir, global_hashes)

        if not items:
            skipped += 1
            print(f"  No items found, skipping")
            # Clean up empty directory
            try:
                os.rmdir(slide_dir)
            except OSError:
                pass
            continue

        # Track which looks each hash appears in
        for item in items:
            h = item["imageHash"]
            if look_id not in global_hashes.get(h, []):
                global_hashes.setdefault(h, []).append(look_id)

        extracted += 1
        print(f"  Extracted {len(items)} items, {len(panel_colors)} panel colors")

        manifest["looks"].append({
            "id": look_id,
            "slideIndex": slide_idx,
            "slideImage": slide_image,
            "category": category,
            "panelColors": panel_colors,
            "items": items,
        })

    # Add shared image info to manifest
    shared_images = {h: looks for h, looks in global_hashes.items() if len(looks) > 1}
    manifest["sharedImages"] = shared_images

    manifest_path = os.path.join(OUT_DIR, "manifest-raw.json")
    with open(manifest_path, "w") as f:
        json.dump(manifest, f, indent=2)

    total_items = sum(len(l["items"]) for l in manifest["looks"])
    print(f"\nDone!")
    print(f"  Extracted: {extracted} looks, {total_items} items")
    print(f"  Skipped: {skipped} slides (title/divider/empty)")
    print(f"  Shared images: {len(shared_images)} images appear in multiple looks")
    print(f"  Manifest: {manifest_path}")


if __name__ == "__main__":
    main()
