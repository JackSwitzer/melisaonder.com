"""
Extract individual items from PPTX slides into a manifest for responsive rendering.
Handles Canva exports where images are embedded as blip fills inside FREEFORM shapes.
Outputs: public/outfits/extracted/<slide-id>/item-<n>.<ext> + manifest.json
"""

import json
import os
import re
import sys
from pptx import Presentation

PPTX_PATH = sys.argv[1] if len(sys.argv) > 1 else os.path.expanduser(
    "~/Downloads/melisaonder.com outfits (1).pptx"
)
OUT_DIR = os.path.join(os.path.dirname(__file__), "..", "public", "outfits", "extracted")
# How many slides to extract (0 = all)
MAX_SLIDES = int(sys.argv[2]) if len(sys.argv) > 2 else 5

EMU_PER_PX = 914400 / 96  # 9525 EMUs per pixel at 96 DPI


def emu_to_px(emu: int) -> int:
    return round(emu / EMU_PER_PX)


def extract_blip_image(shape, slide_part):
    """Extract image blob from a shape's blip fill (Canva freeform pattern)."""
    xml = shape._element.xml
    rids = re.findall(r'r:embed="(rId\d+)"', xml)
    for rid in rids:
        try:
            rel = slide_part.rels[rid]
            if "image" in rel.reltype:
                blob = rel.target_part.blob
                ct = rel.target_part.content_type
                return blob, ct
        except Exception:
            continue
    return None, None


def extract_shapes(shapes, slide_part, out_dir):
    """Extract image shapes — handles standard pictures, freeform blip fills, and groups."""
    items = []

    for shape_idx, shape in enumerate(shapes):
        # Handle grouped shapes — recurse into group
        if shape.shape_type == 6:  # MSO_SHAPE_TYPE.GROUP
            try:
                items.extend(extract_shapes(shape.shapes, slide_part, out_dir))
            except Exception:
                pass
            continue

        img_bytes = None
        content_type = None

        # Method 1: standard picture shape
        if hasattr(shape, "image"):
            try:
                img_bytes = shape.image.blob
                content_type = shape.image.content_type
            except Exception:
                pass

        # Method 2: Canva-style blip fill in freeform/auto shapes
        if img_bytes is None:
            img_bytes, content_type = extract_blip_image(shape, slide_part)

        if img_bytes is None:
            continue

        try:
            ext = content_type.split("/")[-1].replace("jpeg", "jpg")

            left_px = emu_to_px(shape.left)
            top_px = emu_to_px(shape.top)
            w_px = emu_to_px(shape.width)
            h_px = emu_to_px(shape.height)

            fname = f"item-{len(items)}.{ext}"
            fpath = os.path.join(out_dir, fname)
            with open(fpath, "wb") as f:
                f.write(img_bytes)

            rotation = shape.rotation if hasattr(shape, "rotation") else 0

            items.append({
                "file": fname,
                "left": left_px,
                "top": top_px,
                "width": w_px,
                "height": h_px,
                "rotation": rotation,
                "zIndex": len(items),
            })
        except Exception as e:
            print(f"  Skipping shape {shape_idx}: {e}")

    return items


def main():
    print(f"Opening: {PPTX_PATH}")
    prs = Presentation(PPTX_PATH)

    slide_w = emu_to_px(prs.slide_width)
    slide_h = emu_to_px(prs.slide_height)
    print(f"Slide dimensions: {slide_w}x{slide_h} px")

    os.makedirs(OUT_DIR, exist_ok=True)

    manifest = {
        "slideWidth": slide_w,
        "slideHeight": slide_h,
        "looks": [],
    }

    slides = list(prs.slides)
    if MAX_SLIDES > 0:
        slides = slides[:MAX_SLIDES]

    for slide_idx, slide in enumerate(slides):
        slide_id = f"look-{slide_idx + 1}"
        slide_dir = os.path.join(OUT_DIR, slide_id)
        os.makedirs(slide_dir, exist_ok=True)

        print(f"Extracting slide {slide_idx + 1}/{len(slides)}: {slide_id}")

        # Try to get background color
        bg_color = None
        try:
            bg = slide.background
            if bg.fill and bg.fill.fore_color:
                bg_color = f"#{bg.fill.fore_color.rgb}"
        except Exception:
            pass

        items = extract_shapes(slide.shapes, slide.part, slide_dir)
        print(f"  Extracted {len(items)} items")

        manifest["looks"].append({
            "id": slide_id,
            "slideIndex": slide_idx,
            "backgroundColor": bg_color,
            "items": items,
        })

    manifest_path = os.path.join(OUT_DIR, "manifest.json")
    with open(manifest_path, "w") as f:
        json.dump(manifest, f, indent=2)

    print(f"\nDone! Manifest: {manifest_path}")
    print(f"Extracted {sum(len(l['items']) for l in manifest['looks'])} items from {len(manifest['looks'])} slides")


if __name__ == "__main__":
    main()
